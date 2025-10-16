from mcp.server.fastmcp import Context, FastMCP
import sys
import requests
import time
import cv2
import base64
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
import os
import sys
from pathlib import Path
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from datetime import datetime

mcp = FastMCP("Analysis", host="localhost", port=6969)


# video analysis
# -------------------------------------------------------------------------------------------------------------------------------------------------------------------------
def extract_key_frame(video_path, interval_sec=3, max_frames=20):
    """Extract several frames every few seconds."""
    cap = cv2.VideoCapture(video_path)
    if not cap.isOpened():
        raise ValueError(f"Cannot open video file: {video_path}")
    fps = cap.get(cv2.CAP_PROP_FPS)
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    if fps == 0 or total_frames == 0:
        raise ValueError(
            f"Invalid video properties: fps={fps}, total_frames={total_frames}"
        )
    duration = total_frames / fps

    frames_b64 = []
    timestamps = np.arange(0, duration, interval_sec)

    for t in timestamps[:max_frames]:
        cap.set(cv2.CAP_PROP_POS_MSEC, t * 1000)
        ret, frame = cap.read()
        if not ret:
            break
        _, buffer = cv2.imencode(".jpg", frame)
        frames_b64.append(base64.b64encode(buffer).decode("utf-8"))

    cap.release()
    return frames_b64




@mcp.tool()
def analyze_video():
    # Extract the first frame (you can extend this to multiple frames)
    video_path = "uploaded_video.mp4"
    print("extracting....")
    image_b64 = extract_key_frame(video_path)
    print("done")

    # Ollama API endpoint
    url = "http://localhost:11434/api/chat"
    question = "summarize this image"

    # Prepare request payload
    payload = {
        "model": "minicpm-v",
        "messages": [
            {
                "role": "user",
                "content": question,
                "images": image_b64,
            },
        ],
        "stream": False,
    }

    # Send request
    response = requests.post(url, json=payload, stream=False)
    if response.status_code != 200:
        raise RuntimeError(f"API error: {response.text}")
    
    result = response.json()["message"]["content"]
    
    return result


# chat history summary
# -------------------------------------------------------------------------------------------------------------------------------------------------------------------------
def save_to_pdf(text, output_path):
    """Save text content to a simple PDF file."""
    # Ensure directory exists
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)

    c = canvas.Canvas(output_path, pagesize=A4)
    width, height = A4
    text_object = c.beginText(40, height - 50)
    text_object.setFont("Helvetica", 12)

    # Wrap long lines manually
    max_width = 90  # characters per line
    for line in text.splitlines():
        for wrapped in [line[i:i + max_width] for i in range(0, len(line), max_width)]:
            text_object.textLine(wrapped)

    c.drawText(text_object)
    c.save()

@mcp.tool()
def summarize_chat_history():
    """
    Reads chat history from a text file and sends it to Ollama for summarization. Always respond that summary is saved in the returned pdf path
    """
    file_path = "chat_history.txt"
    # Step 1: Read file content
    with open(file_path, "r", encoding="utf-8") as f:
        chat_text = f.read()

    # Step 2: Define Ollama API endpoint
    url = "http://localhost:11434/api/generate"  # default Ollama API endpoint

    # Step 3: Define the request payload
    payload = {
        "model": "qwen3:8b",
        "prompt": (
            "Summarize the following discussion between a user and an AI assistant. "
            "Focus on key insights, sentiment, and conclusions:\n\n"
            f"{chat_text}"
        ),
        "stream": False,  # disable streaming for simplicity
    }

    # Step 4: Send HTTP POST request
    try:
        response = requests.post(url, json=payload)
        response.raise_for_status()
    except requests.exceptions.RequestException as e:
        print("Failed to connect to Ollama:", e)
        return None

    # Step 5: Extract the generated summary
    data = response.json()
    summary = data.get("response", "").strip()

    # Save result to PDF
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    pdf_path = f"summary_{timestamp}.pdf"
    save_to_pdf(summary, pdf_path)

    return summary,pdf_path


# powerpoint generation
# -------------------------------------------------------------------------------------------------------------------------------------------------------------------------
PPT_PATH = "current.pptx"


def get_presentation():
    """Load or create a shared presentation"""
    if os.path.exists(PPT_PATH):
        return Presentation(PPT_PATH)
    else:
        prs = Presentation()
        prs.save(PPT_PATH)
        return prs


@mcp.tool()
def create_ppt():
    """
    Creates a new PowerPoint presentation file in the current directory. Use this before adding any slides.
    """
    prs = Presentation()
    prs.save(PPT_PATH)
    return "Created a new presentation."


@mcp.tool()
def add_slide(layout: int) -> int:
    """
    Adds a new slide to the current presentation using the specified layout index (e.g., layout=0 for title slide, 1 for title + content).

    Args:
    layout (int): layout index
    """
    prs = get_presentation()
    slide_layout = prs.slide_layouts[layout]
    prs.slides.add_slide(slide_layout)
    prs.save(PPT_PATH)
    return f"Added a new slide using layout {layout}."


@mcp.tool()
def write_text(text: str, slide_index: int, placeholder: int) -> str:
    """
    Writes or updates text in the given placeholder on a specific slide.

    Args:
    text (str): Content or title need to be written
    slide_index (int): the slide number
    placeholder (int): the number of the text box, eg. 0 for title and 1 for content box
    """
    prs = get_presentation()
    slide = prs.slides[slide_index]

    if placeholder < len(slide.placeholders):
        slide.placeholders[placeholder].text = text
        for paragraph in slide.placeholders[placeholder].text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(25)
    else:
        # Add a new text box
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = text

        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(15)

    prs.save(PPT_PATH)
    return f"Added text to slide {slide_index}."


@mcp.tool()
def save_ppt(filename):
    """
    Rename the ppt name created based on filename.

    Args:
    filename (str): filename to be renamed
    """
    if os.path.exists(PPT_PATH):
        os.rename(PPT_PATH, f"{filename}.pptx")
        return f"Saved presentation as {filename}"
    else:
        return "No presentation found to save."


# generate ppt template
# -------------------------------------------------------------------------------------------------------------------------------------------------------------------------

# Configuration
# OLLAMA_HOST = "http://localhost:11434"  # change if running on another host/IP
# MODEL = "qwen3:8b"  # or whatever model you want (e.g., llama3, mistral, etc.)


# def read_transcript(file_path: str) -> str:
#     path = Path(file_path)
#     if not path.exists():
#         print(f"File not found: {file_path}")
#         sys.exit(1)
#     return path.read_text(encoding="utf-8").strip()


# def save_to_file(content: str) -> str:
#     output_file = "ppt_template.txt"
#     Path(output_file).write_text(content, encoding="utf-8")
#     return output_file


# @mcp.tool()
# def generate_ppt_template():
#     transcript = read_transcript("transcript.txt")

#     # Prompt template
#     prompt = f"""
#     You are an AI assistant that converts a transcript into PowerPoint slide templates.

#     Rules:
#     - Each slide begins with "slide X:"
#     - Slide 1 = title only
#     - Slide 2 = title and content
#     - Last slide = summary

#     Convert the following transcript into slides:

#     {transcript}
#     """

#     # Send request to Ollama API
#     response = requests.post(
#         f"{OLLAMA_HOST}/api/generate",
#         json={"model": MODEL, "prompt": prompt, "stream": False},
#     )

#     if response.status_code != 200:
#         raise RuntimeError(f"Ollama error {response.status_code}: {response.text}")

#     data = response.json()
#     template = data.get("response", "").strip()
#     output_path = save_to_file(template)
#     return template, output_path


if __name__ == "__main__":
    try:
        # Run the server
        mcp.run(transport="streamable-http")

    except KeyboardInterrupt:
        print("Server shutting down gracefully...")
        print("Server has been shut down.")

    except Exception as e:
        print(f"Unexpected error: {e}")
        sys.exit(1)

    finally:
        print("Thank you for using MCP Server!")
